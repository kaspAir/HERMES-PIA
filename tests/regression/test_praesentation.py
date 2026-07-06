"""PIA-Upload + Präsentationsgenerierung: Parser, Builder, Vorlagen-Vorrang, Routen."""
import base64
import io

import docx
import pytest
from pptx import Presentation

from app.config import Config
from app.domains.praesentation.parser import parse_pia
from app.domains.praesentation.service import PraesentationService
from app.factory import create_app


# ---- Hilfen: Beispiel-PIA (.docx) und Vorlage (.pptx) im Speicher ------ #

LANGES_RISIKO = ("Schlüsselpersonen aus dem IT-Betrieb sind wegen operativer Belastung "
                 "nicht ausreichend für die Erarbeitung der Studie, der Schutzbedarfsanalyse "
                 "und des Projektmanagementplans verfügbar, wodurch die fristgerechte "
                 "Fertigstellung der Initialisierungsergebnisse gefährdet wird.")


def _beispiel_pia_bytes():
    d = docx.Document()
    # Kopftabelle (vor der ersten Überschrift) -> Metadaten
    meta = d.add_table(rows=2, cols=2)
    meta.cell(0, 0).text = "Projektleiter/in"
    meta.cell(0, 1).text = "Petra Muster"
    meta.cell(1, 0).text = "Auftraggeber/in"
    meta.cell(1, 1).text = "Hans Beispiel"

    d.add_heading("Ausgangslage", level=1)
    d.add_paragraph("Das System X wird abgelöst. Die Betriebssicherheit muss steigen. "
                    "Alle Einheiten sind betroffen.")

    d.add_heading("Ziele der Phase Initialisierung", level=2)
    t = d.add_table(rows=2, cols=5)
    for i, h in enumerate(("Nr.", "Kategorie", "Beschreibung", "Messgrösse", "Priorität")):
        t.cell(0, i).text = h
    t.cell(1, 0).text = "01"
    t.cell(1, 2).text = "Die Studie liegt als Entscheidungsgrundlage vor."
    t.cell(1, 3).text = "Studie abgenommen"

    d.add_heading("Personalaufwand", level=2)
    t = d.add_table(rows=3, cols=3)
    for i, h in enumerate(("Rolle", "Name", "Aufwand in PT")):
        t.cell(0, i).text = h
    t.cell(1, 0).text = "Projektleiterin"; t.cell(1, 1).text = "Petra Muster"
    t.cell(1, 2).text = "25"
    t.cell(2, 0).text = "Externe Fachexpertise (extern)"; t.cell(2, 2).text = "20"

    d.add_heading("Kosten (in CHF inkl. MwSt.)", level=2)
    t = d.add_table(rows=5, cols=2)
    t.cell(0, 0).text = "Phase"; t.cell(0, 1).text = "Betrag"
    t.cell(1, 0).text = "Interne Personalkosten"; t.cell(1, 1).text = "50400"
    t.cell(2, 0).text = "Externe Fachexpertise (extern)"; t.cell(2, 1).text = "36000"
    t.cell(3, 0).text = "Summe externe Kosten"; t.cell(3, 1).text = "36000"
    t.cell(4, 0).text = "Total Initialisierung"; t.cell(4, 1).text = "86400"

    d.add_heading("Ergebnisse und Termine", level=2)
    t = d.add_table(rows=4, cols=5)
    for i, h in enumerate(("Nr.", "Lieferergebnisse (abnahmerelevant)", "Liefertermin",
                           "Abnahme durch (Rolle)", "Prüfmethode")):
        t.cell(0, i).text = h
    t.cell(1, 1).text = "Rechtsgrundlagenanalyse"; t.cell(1, 2).text = "05.10.2026"
    t.cell(1, 3).text = "Projektleiter"
    t.cell(2, 1).text = "Studie"; t.cell(2, 2).text = "26.10.2026"; t.cell(2, 3).text = "Projektleiter"
    t.cell(3, 1).text = "Meilenstein Durchführungsfreigabe"; t.cell(3, 2).text = "14.12.2026"
    t.cell(3, 3).text = "Auftraggeber"

    d.add_heading("Risiken", level=1)
    t = d.add_table(rows=3, cols=8)
    for i, h in enumerate(("Nr.", "Risikobeschreibung", "Eintrittswahrscheinlichkeit",
                           "Auswirkungsgrad", "Risikozahl", "Massnahmen",
                           "Verantwortung", "Termin")):
        t.cell(0, i).text = h
    t.cell(1, 1).text = LANGES_RISIKO
    t.cell(1, 2).text = "Hoch"; t.cell(1, 3).text = "Mittel"
    t.cell(2, 1).text = "Externe Expertise fehlt."
    t.cell(2, 4).text = "9"                     # EW/AG leer -> Fallback über Risikozahl

    d.add_heading("Dokument-Protokoll", level=1)  # Ende-Marker
    d.add_paragraph("Wird ignoriert.")

    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _leere_pptx_bytes(width=None):
    prs = Presentation()
    if width:
        prs.slide_width = width
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---- Parser ------------------------------------------------------------ #

def test_parser_liest_abschnitte_und_tabellen():
    pia = parse_pia(_beispiel_pia_bytes())
    assert pia["projektleiter"] == "Petra Muster"
    assert pia["auftraggeber"] == "Hans Beispiel"
    assert "Betriebssicherheit" in pia["ausgangslage"]
    assert pia["ziele"][0]["beschreibung"].startswith("Die Studie")
    assert {p["rolle"]: p["aufwand"] for p in pia["personalaufwand"]} == {
        "Projektleiterin": 25, "Externe Fachexpertise (extern)": 20}
    assert pia["personalaufwand"][0]["name"] == "Petra Muster"
    assert any(k["position"] == "Total Initialisierung" and k["betrag"] == 86400
               for k in pia["kosten"])
    assert pia["termine"][0]["ergebnis"] == "Rechtsgrundlagenanalyse"
    # "Abnahme durch (Rolle)" darf NICHT auf "(abnahmerelevant)" matchen:
    assert pia["termine"][0]["abnahme"] == "Projektleiter"
    assert pia["termine"][2]["abnahme"] == "Auftraggeber"
    # Risiko 1: EW/AG aus Text; Risiko 2: Fallback aus Risikozahl 9 -> (3,3)
    assert (pia["risiken"][0]["ew"], pia["risiken"][0]["ag"]) == (3, 2)
    assert (pia["risiken"][1]["ew"], pia["risiken"][1]["ag"]) == (3, 3)


# ---- Builder ------------------------------------------------------------ #

def test_builder_erzeugt_praesentation_ohne_vorlage():
    svc = PraesentationService(llm=None)
    buf = svc.generate_from_docx(_beispiel_pia_bytes(), template_bytes=None,
                                 fallback_name="Fallback", datum="05.07.2026")
    prs = Presentation(buf)
    # Titel, Ausgangslage, Ziele, Termine, Personal, Kosten, Matrix, Risiken, Antrag
    assert len(prs.slides) >= 8
    titel_texte = " ".join(s.text for s in prs.slides[0].shapes if s.has_text_frame
                           for s in [s])
    # Projektname kommt aus dem Dokument-Fallback (erste Absätze fehlen -> fallback_name)
    assert "Projektinitialisierungsauftrag" in _slide_text(prs.slides[0])


def _slide_text(slide):
    return " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def test_builder_baut_auf_vorlage_auf():
    breite = 12192000  # 16:9-Breite in EMU
    vorlage = _leere_pptx_bytes(width=breite)
    svc = PraesentationService(llm=None)
    buf = svc.generate_from_docx(_beispiel_pia_bytes(), template_bytes=vorlage,
                                 fallback_name="X", datum="")
    prs = Presentation(buf)
    assert prs.slide_width == breite            # Foliengrösse der Vorlage übernommen
    assert len(prs.slides) >= 8


def _vorlage_mit_beispielfolien():
    """Vorlage wie eine Firmenvorlage: 3 Beispiel-Folien auf unterschiedlichen
    Layouts (erste = Titeldesign, letzte = Schlussdesign)."""
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "BEISPIEL-MARKER-1"
    prs.slides.add_slide(prs.slide_layouts[1])
    s3 = prs.slides.add_slide(prs.slide_layouts[2])
    s3.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "BEISPIEL-MARKER-3"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), prs.slide_layouts[0].name, prs.slide_layouts[2].name


def test_beispielfolien_raus_titel_und_schluss_aus_vorlage():
    vorlage, titel_layout, schluss_layout = _vorlage_mit_beispielfolien()
    svc = PraesentationService(llm=None)
    buf = svc.generate_from_docx(_beispiel_pia_bytes(), template_bytes=vorlage,
                                 fallback_name="X", datum="06.07.2026")
    prs = Presentation(buf)
    alle_texte = " ".join(_slide_text(s) for s in prs.slides)
    assert "BEISPIEL-MARKER" not in alle_texte            # Beispiel-Folien entfernt
    assert prs.slides[0].slide_layout.name == titel_layout   # Titeldesign der Vorlage
    assert prs.slides[-1].slide_layout.name == schluss_layout  # Schlussdesign (z.B. blau)
    assert "Besten Dank" in _slide_text(prs.slides[-1])


def _prst_geom(shape):
    el = shape._element.find(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
    return el.get("prst") if el is not None else None


def test_gantt_folie_mit_balken_und_meilenstein_rauten():
    svc = PraesentationService(llm=None)
    buf = svc.generate_from_docx(_beispiel_pia_bytes(), None, "X", "")
    prs = Presentation(buf)
    gantt = next((s for s in prs.slides if "Terminplan" in _slide_text(s)), None)
    assert gantt is not None
    formen = [_prst_geom(sh) for sh in gantt.shapes]
    assert formen.count("diamond") == 1        # 1 Meilenstein -> 1 Raute (Dauer 0)
    # 2 Ergebnis-Balken + 1 helle Verlängerung (Studie -> Meilenstein, lückenlos)
    assert formen.count("roundRect") == 3
    # Echte Daten an Balken/Rauten (statt Kalendertage ab Start)
    text = _slide_text(gantt)
    assert "26.10.26" in text and "14.12.26" in text
    assert "Kalendertage" not in text


def test_personalaufwand_namen_rechtsbuendig_mit_nn_platzhalter():
    from pptx.enum.text import PP_ALIGN
    svc = PraesentationService(llm=None)
    buf = svc.generate_from_docx(_beispiel_pia_bytes(), None, "X", "")
    prs = Presentation(buf)
    folie = next(s for s in prs.slides if "Personalaufwand" in _slide_text(s))
    text = _slide_text(folie)
    assert "Petra Muster" in text               # erfasster Name erscheint
    assert "N. N." in text                      # fehlender Name -> Platzhalter
    assert "25 PT" in text and "20 PT" in text  # Werte an den Balken
    # Namenszeile ist rechtsbündig
    for sh in folie.shapes:
        if sh.has_text_frame and "Petra Muster" in sh.text_frame.text:
            name_absatz = sh.text_frame.paragraphs[1]
            assert name_absatz.alignment == PP_ALIGN.RIGHT


def test_projektplan_msproject_und_excel():
    from app.domains.praesentation.parser import parse_pia
    from app.domains.praesentation import projektplan
    import xml.etree.ElementTree as ET

    pia = parse_pia(_beispiel_pia_bytes())
    eintraege = projektplan.plan_eintraege(pia["termine"])
    assert len(eintraege) == 3
    # Kaskade: Studie startet am Termin der Rechtsgrundlagenanalyse
    namen = [e[0] for e in eintraege]
    studie = eintraege[namen.index("Studie")]
    assert studie[1].strftime("%d.%m.%Y") == "05.10.2026"
    assert studie[2].strftime("%d.%m.%Y") == "26.10.2026"
    ms = next(e for e in eintraege if e[3])
    assert ms[1] == ms[2]                       # Meilenstein: Dauer 0

    xml_bytes = projektplan.build_msproject_xml(eintraege, "P Demo")
    root = ET.fromstring(xml_bytes)
    ns = "{http://schemas.microsoft.com/project}"
    tasks = root.findall(f"{ns}Tasks/{ns}Task")
    assert len(tasks) == 3
    flags = [t.findtext(f"{ns}Milestone") for t in tasks]
    assert flags.count("1") == 1
    # Project prueft die xsd:sequence STRIKT - Element-Reihenfolge muss dem
    # mspdi-Schema entsprechen (2010er-Felder Active/Manual/... am Ende).
    for task in tasks:
        kinder = tuple(el.tag.replace(ns, "") for el in task)
        assert kinder == projektplan._TASK_ELEMENT_ORDER

    excel = projektplan.build_excel(eintraege, "P Demo")
    assert excel.startswith(b"PK")              # gültige xlsx-(ZIP-)Datei
    # Gantt-Raster: Monatsköpfe, Meilenstein-Raute und Balkenfarbe vorhanden
    import zipfile
    with zipfile.ZipFile(io.BytesIO(excel)) as z:
        strings = z.read("xl/sharedStrings.xml").decode("utf-8")
        styles = z.read("xl/styles.xml").decode("utf-8")
    assert "Okt 26" in strings and "Dez 26" in strings   # Monatsblöcke über den KW-Spalten
    assert "◆" in strings                                # Meilenstein-Raute im Raster
    assert "2E75B6" in styles                            # Balken-Füllfarbe


def test_projektplan_routen(app):
    c, org_id, pid, eid = _client_mit_projekt(app)
    # Ohne hochgeladenen PIA: klarer Hinweis
    assert c.get(f"/projekt/{pid}/ergebnis/{eid}/projektplan/excel/x.xlsx").status_code == 400
    c.post(f"/projekt/{pid}/ergebnis/{eid}/dokument",
           json={"filename": "PIA.docx", "data": _b64(_beispiel_pia_bytes())})
    r = c.get(f"/projekt/{pid}/ergebnis/{eid}/projektplan/excel/20260706_P_Demo_Projektplan.xlsx")
    assert r.status_code == 200 and r.data.startswith(b"PK")
    assert "20260706_P_Demo_Projektplan.xlsx" in r.headers.get("Content-Disposition", "")
    r = c.get(f"/projekt/{pid}/ergebnis/{eid}/projektplan/msproject/20260706_P_Demo_Projektplan.xml")
    assert r.status_code == 200 and b"schemas.microsoft.com/project" in r.data
    assert c.get(f"/projekt/{pid}/ergebnis/{eid}/projektplan/quatsch/x.xml").status_code == 404


def test_risiken_volltext_notizen_und_keine_auslassungspunkte():
    svc = PraesentationService(llm=None)
    buf = svc.generate_from_docx(_beispiel_pia_bytes(), None, "X", "")
    prs = Presentation(buf)
    # Volle Risikobeschreibung in einer Tabelle - keine Kürzung.
    zelltexte = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if getattr(sh, "has_table", False):
                zelltexte += [c.text for r in sh.table.rows for c in r.cells]
    assert any(LANGES_RISIKO == t for t in zelltexte)
    # HERMES-Terminologie: "Ergebnis", nicht "Lieferergebnis".
    assert "Ergebnis" in zelltexte and not any("Lieferergebnis" in t for t in zelltexte)
    # NIRGENDS abgeschnittene Sätze mit '…' - weder auf Folien noch in Notizen.
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                assert "…" not in sh.text_frame.text
        assert "…" not in " ".join(zelltexte)
        assert slide.has_notes_slide
        notizen = slide.notes_slide.notes_text_frame.text
        assert notizen.strip() and "…" not in notizen and "Lieferergebnis" not in notizen


def test_notizen_auch_ohne_notizen_platzhalter():
    """Firmenvorlagen ohne Body-Platzhalter im Notizen-Master (notes_text_frame
    = None) dürfen die Generierung nicht crashen – der Platzhalter wird als
    XML nachgerüstet und die Notiz gesetzt."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    from app.domains.praesentation.service import _Builder

    prs = Presentation()
    builder = _Builder(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    notes = slide.notes_slide
    # Situation der Firmenvorlage simulieren: Body-Platzhalter entfernen.
    for ph in list(notes.placeholders):
        if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
            ph._element.getparent().remove(ph._element)
    assert notes.notes_text_frame is None

    builder._notizen(slide, "Sprechnotiz für den Projektleiter.")
    assert slide.notes_slide.notes_text_frame is not None
    assert "Sprechnotiz" in slide.notes_slide.notes_text_frame.text


def test_parser_risiken_mit_verschobenen_spalten():
    """Datenzeilen kuerzer als der Kopf (leere SDT-Zellen) -> Felder inhaltlich finden."""
    from app.domains.praesentation.parser import _parse_risiken
    rows = [
        ["Nr.", "Risikobeschreibung", "Eintrittswahrscheinlichkeit", "Auswirkungsgrad",
         "Risikozahl", "Massnahmen", "Verantwortung", "Termin"],
        ["01", "Externe Fachexpertise ist nicht rechtzeitig verfuegbar und gefaehrdet die Studie.",
         "9", "Fruehzeitig Marktsondierung einleiten und Anbieter anfragen.",
         "Projektleiter", "laufend"],
    ]
    out = _parse_risiken(rows)
    assert len(out) == 1
    assert (out[0]["ew"], out[0]["ag"]) == (3, 3)          # Risikozahl 9 -> hoch/hoch
    assert out[0]["beschreibung"].startswith("Externe Fachexpertise")


# ---- App-Fixture für Service/Routen ------------------------------------- #

@pytest.fixture
def app(tmp_path):
    from app.shared.database import SessionLocal
    db_path = str(tmp_path / "praes.db").replace("\\", "/")

    class _Cfg(Config):
        DATABASE_URL = "sqlite:///" + db_path
        SUPERADMIN_EMAIL = "betreiber@test.ch"
        SUPERADMIN_PASSWORD = "pw-super"
        SECRET_KEY = "test-secret"

    SessionLocal.remove()
    application = create_app(_Cfg)
    SessionLocal.remove()
    yield application
    SessionLocal.remove()


def _client_mit_projekt(app):
    auth = app.auth_service
    org_id = auth.create_org("Org").id      # ID sofort sichern (Detached-Falle)
    auth.create_user("pl@org.ch", "pw", org_id=org_id,
                     can_read=True, can_write=True, can_delete=True)
    c = app.test_client()
    c.post("/login", data={"email": "pl@org.ch", "password": "pw"})
    c.post("/interview/start", data={"project_name": "P Demo", "projektleiter": "PL"})
    pid = app.projekt_service.projekte_for_org(org_id)[0].id
    eid = app.projekt_service.ergebnisse(pid)[0].id
    return c, org_id, pid, eid


# ---- Vorlagen-Vorrang ----------------------------------------------------- #

def test_vorlage_projekt_schlaegt_org(app):
    svc = app.projekt_service
    c, org_id, pid, eid = _client_mit_projekt(app)
    projekt = svc.get_projekt(pid)
    assert svc.resolve_vorlage(projekt) is None
    svc.add_vorlage("org.pptx", b"PKorg", org_id=org_id, projekt_id=None)
    assert svc.resolve_vorlage(projekt).filename == "org.pptx"
    svc.add_vorlage("projekt.pptx", b"PKprj", org_id=org_id, projekt_id=pid)
    assert svc.resolve_vorlage(projekt).filename == "projekt.pptx"


def test_pmo_seite_und_orgweite_vorlage(app):
    """PMO-Bereich: organisationsweite Vorlage hochladen; Projekt-Upload bleibt
    projektspezifisch und übersteuert die PMO-Vorlage."""
    c, org_id, pid, eid = _client_mit_projekt(app)

    r = c.get("/pmo")
    assert r.status_code == 200
    html = r.get_data(as_text=True)
    assert "PMO" in html and "Noch keine Vorlage" in html

    # PMO-Vorlage hochladen -> gilt organisationsweit
    r = c.post("/pmo/vorlage",
               json={"filename": "pmo.pptx", "data": _b64(_leere_pptx_bytes())})
    assert r.status_code == 200
    svc = app.projekt_service
    assert svc.org_vorlage(org_id).filename == "pmo.pptx"
    assert svc.resolve_vorlage(svc.get_projekt(pid)).filename == "pmo.pptx"
    assert "pmo.pptx" in c.get("/pmo").get_data(as_text=True)

    # Projekt-Upload ist IMMER projektspezifisch (auch mit scope-Altparameter)
    c.post(f"/projekt/{pid}/vorlage",
           json={"filename": "override.pptx", "data": _b64(_leere_pptx_bytes()),
                 "scope": "org"})
    assert svc.resolve_vorlage(svc.get_projekt(pid)).filename == "override.pptx"
    assert svc.org_vorlage(org_id).filename == "pmo.pptx"   # PMO-Vorlage unangetastet


# ---- Routen ----------------------------------------------------------------- #

def _b64(data):
    return base64.b64encode(data).decode()


def test_upload_download_und_status(app):
    c, org_id, pid, eid = _client_mit_projekt(app)
    pia = _beispiel_pia_bytes()
    r = c.post(f"/projekt/{pid}/ergebnis/{eid}/dokument",
               json={"filename": "PIA_v0.9.docx", "data": _b64(pia)})
    assert r.status_code == 200
    dok = app.projekt_service.latest_dokument(eid, art="freigabe")
    assert dok.filename == "PIA_v0.9.docx" and dok.size == len(pia)
    # Status des Ergebnisses wechselt auf "zur Freigabe"
    erg = app.projekt_service.ergebnisse(pid)[0]
    assert erg.status == "zur Freigabe"
    # Download liefert die Originaldatei (Dateiname in der URL, wegen Proxy)
    r = c.get(f"/projekt/{pid}/dokument/{dok.id}/{dok.filename}")
    assert r.status_code == 200 and r.data == pia


def test_upload_validierung(app):
    c, org_id, pid, eid = _client_mit_projekt(app)
    url = f"/projekt/{pid}/ergebnis/{eid}/dokument"
    assert c.post(url, json={"filename": "x.txt", "data": _b64(b"PKx")}).status_code == 400
    assert c.post(url, json={"filename": "x.docx", "data": ""}).status_code == 400
    assert c.post(url, json={"filename": "x.docx", "data": _b64(b"kein zip")}).status_code == 400


def test_praesentation_route(app):
    c, org_id, pid, eid = _client_mit_projekt(app)
    # Ohne hochgeladenen PIA: klarer Hinweis statt Absturz
    assert c.get(f"/projekt/{pid}/ergebnis/{eid}/praesentation/x.pptx").status_code == 400
    c.post(f"/projekt/{pid}/ergebnis/{eid}/dokument",
           json={"filename": "PIA.docx", "data": _b64(_beispiel_pia_bytes())})
    # Mit Projekt-Vorlage
    c.post(f"/projekt/{pid}/vorlage",
           json={"filename": "v.pptx", "data": _b64(_leere_pptx_bytes()), "scope": "projekt"})
    # Dateiname steht IN DER URL (Proxy verschluckt Content-Disposition) und
    # wird zusätzlich als download_name gesetzt.
    r = c.get(f"/projekt/{pid}/ergebnis/{eid}/praesentation/20260706_P_Demo.pptx")
    assert r.status_code == 200
    assert "presentationml" in r.mimetype
    assert "20260706_P_Demo.pptx" in r.headers.get("Content-Disposition", "")
    prs = Presentation(io.BytesIO(r.data))
    assert len(prs.slides) >= 8
    # Die Projektseite verlinkt mit yyyymmdd_Projektname.pptx in der URL.
    import re as _re
    html = c.get(f"/projekt/{pid}").get_data(as_text=True)
    assert _re.search(r"/praesentation/\d{8}_P_Demo\.pptx", html)
    assert _re.search(r"/projektplan/msproject/\d{8}_P_Demo_Projektplan\.xml", html)
    assert _re.search(r"/projektplan/excel/\d{8}_P_Demo_Projektplan\.xlsx", html)


def test_dokumente_fremder_org_gesperrt(app):
    c, org_id, pid, eid = _client_mit_projekt(app)
    pia = _beispiel_pia_bytes()
    c.post(f"/projekt/{pid}/ergebnis/{eid}/dokument",
           json={"filename": "PIA.docx", "data": _b64(pia)})
    dok_id = app.projekt_service.latest_dokument(eid).id

    auth = app.auth_service
    other_id = auth.create_org("Andere").id
    auth.create_user("fremd@x.ch", "pw", org_id=other_id,
                     can_read=True, can_write=True, can_delete=False)
    cb = app.test_client()
    cb.post("/login", data={"email": "fremd@x.ch", "password": "pw"})
    assert cb.get(f"/projekt/{pid}/dokument/{dok_id}/PIA.docx").status_code == 403
    assert cb.post(f"/projekt/{pid}/ergebnis/{eid}/dokument",
                   json={"filename": "x.docx", "data": _b64(pia)}).status_code == 403


def test_delete_projekt_raeumt_dokumente_und_vorlagen(app):
    c, org_id, pid, eid = _client_mit_projekt(app)
    c.post(f"/projekt/{pid}/ergebnis/{eid}/dokument",
           json={"filename": "PIA.docx", "data": _b64(_beispiel_pia_bytes())})
    c.post(f"/projekt/{pid}/vorlage",
           json={"filename": "v.pptx", "data": _b64(_leere_pptx_bytes()), "scope": "projekt"})
    dok_id = app.projekt_service.latest_dokument(eid).id
    c.post(f"/projekt/{pid}/delete")
    assert app.projekt_service.get_projekt(pid) is None
    assert app.projekt_service.get_dokument(dok_id) is None
