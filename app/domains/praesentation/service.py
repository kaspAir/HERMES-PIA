"""Erzeugt die Präsentation für Auftraggeber/Projektausschuss aus dem
hochgeladenen (freigabebereiten) PIA.

Vorlagenlogik: Die hochgeladene .pptx liefert Design UND Layout-Zuordnung.
* Beispiel-Folien der Vorlage werden ENTFERNT – vorher merkt sich der Builder
  ihre Layouts: das Layout der ersten Folie wird zum Titel-Layout, das der
  letzten zur Schlussfolie (z.B. das blaue Dank-Deckblatt einer Firmenvorlage).
* Inhaltsfolien nutzen ein Layout MIT Titel-Platzhalter (z.B. «Nur Titel») –
  damit bleiben Logo, Design und Titelposition der Vorlage erhalten.
* Fusszeilen-/Datums-/Seitenzahl-Platzhalter des Layouts werden in jede Folie
  übernommen (python-pptx klont sie sonst nicht -> Fusszeile fehlte).

Stilregeln: KEINE mit «…» abgeschnittenen Sätze (Volltext in Tabellen),
schwarze Schrift in allen erzeugten Tabellen, HERMES-Terminologie («Ergebnis»,
nicht «Lieferergebnis»). Meilensteine im Gantt als Raute mit Dauer 0.
Jede Folie erhält Notizen mit Sprechhinweisen für den Projektleiter.
"""
import io
import json
import re
from copy import deepcopy
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.util import Emu, Pt

from app.domains.praesentation.parser import parse_pia

# Ampelfarben der Risikomatrix (dezent, an Excel-Klassiker angelehnt).
_GRUEN = RGBColor(0xC6, 0xEF, 0xCE)
_GELB = RGBColor(0xFF, 0xEB, 0x9C)
_ROT = RGBColor(0xFF, 0xC7, 0xCE)
_GRAU = RGBColor(0xF2, 0xF2, 0xF2)
_SCHWARZ = RGBColor(0x00, 0x00, 0x00)
_BALKENBLAU = RGBColor(0x2E, 0x75, 0xB6)
_RAUTENDUNKEL = RGBColor(0x1F, 0x38, 0x64)
_LINIENGRAU = RGBColor(0xD9, 0xD9, 0xD9)

_TITEL_TYPEN = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
_DEKOR_TYPEN = (PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.SLIDE_NUMBER)
_STUFEN_TEXT = {1: "Tief", 2: "Mittel", 3: "Hoch"}
_MONATE = ("Jan", "Feb", "Mär", "Apr", "Mai", "Jun",
           "Jul", "Aug", "Sep", "Okt", "Nov", "Dez")


class PraesentationService:

    def __init__(self, llm=None):
        self.llm = llm

    def generate_from_docx(self, pia_bytes, template_bytes=None,
                           fallback_name="Projekt", datum=""):
        """PIA-.docx + optionale .pptx-Vorlage -> Präsentation (BytesIO)."""
        pia = parse_pia(pia_bytes)
        if not pia.get("projektname"):
            pia["projektname"] = fallback_name

        prs = Presentation(io.BytesIO(template_bytes)) if template_bytes else Presentation()
        b = _Builder(prs)

        b.titel(pia, datum)
        b.ausgangslage(pia, self._ausgangslage_bullets(pia))
        b.ziele(pia)
        b.termine(pia)
        b.gantt(pia)
        b.personalaufwand(pia)
        b.kosten(pia)
        b.risiken(pia)
        b.antrag(pia)
        b.schluss(pia, datum)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    # ---- Zusammenfassung (LLM mit deterministischem Fallback) --------- #

    def _ausgangslage_bullets(self, pia):
        text = (pia.get("ausgangslage") or "").strip()
        if not text:
            return []
        if self.llm:
            try:
                raw = self.llm.complete(
                    "Du fasst die Ausgangslage eines Projektinitialisierungsauftrags "
                    "fuer eine Praesentation vor Auftraggeber und Projektausschuss "
                    "zusammen. Sachlicher Behoerdenstil, nur VOLLSTAENDIGE Saetze oder "
                    "Stichpunkte (nie mit '...' abbrechen), keine Erfindungen - nur, "
                    "was im Text steht. Antworte ausschliesslich mit validem JSON.",
                    [{"role": "user", "content":
                        f"Ausgangslage:\n{text}\n\n"
                        'Gib 3-5 praegnante Stichpunkte (je max. 15 Woerter) als '
                        'JSON zurueck: {"bullets": ["...", "..."]}'}],
                    max_tokens=600,
                )
                m = re.search(r"\{.*\}", raw, re.S)
                data = json.loads(m.group()) if m else {}
                bullets = [str(x).strip() for x in data.get("bullets", []) if str(x).strip()]
                if bullets:
                    return bullets[:5]
            except Exception:  # noqa: BLE001 – Fallback übernimmt
                pass
        haupt = text.split("Komplexitätseinschätzung")[0]
        saetze = [s.strip() for s in re.split(r"(?<=[.!?])\s+", haupt) if s.strip()]
        return saetze[:4]


# ---------------------------------------------------------------------- #
# Hilfen                                                                   #
# ---------------------------------------------------------------------- #

def _alle_layouts(prs):
    out = []
    for master in prs.slide_masters:
        out.extend(master.slide_layouts)
    return out


def _titel_platzhalter(shape_container):
    for ph in shape_container.placeholders:
        if ph.placeholder_format.type in _TITEL_TYPEN:
            return ph
    return None


def _platzhalter_nach_typ(slide, *typen):
    for ph in slide.placeholders:
        if ph.placeholder_format.type in typen:
            return ph
    return None


def _parse_datum(text):
    m = re.search(r"\d{2}\.\d{2}\.\d{4}", str(text or ""))
    if not m:
        return None
    try:
        return datetime.strptime(m.group(), "%d.%m.%Y")
    except ValueError:
        return None


def _ohne_klammern(text):
    """Saubere Kurzform ohne Abschneiden: Klammerzusätze entfernen."""
    return re.sub(r"\s*\([^)]*\)", "", str(text or "")).strip()


def _erster_satz(text):
    teile = re.split(r"(?<=[.!?])\s+", str(text or "").strip())
    return teile[0] if teile and teile[0] else str(text or "").strip()


def _monatsanfaenge(start, ende):
    d = datetime(start.year, start.month, 1)
    out = []
    while d <= ende:
        if d >= start:
            out.append(d)
        d = datetime(d.year + (1 if d.month == 12 else 0), d.month % 12 + 1, 1)
    return out


def _ist_summe(position):
    p = (position or "").lower()
    return "summe" in p or "total" in p


class _Builder:

    def __init__(self, prs):
        self.prs = prs
        self.sw = prs.slide_width
        self.sh = prs.slide_height
        vorlage_layouts = [s.slide_layout for s in prs.slides]
        self._remove_all_slides(prs)

        self.content_layout = self._pick_content_layout(prs)
        self.title_layout = (vorlage_layouts[0] if vorlage_layouts
                             else self._pick_title_layout(prs))
        self.closing_layout = (vorlage_layouts[-1]
                               if len(vorlage_layouts) >= 2 else None)

    # ---- Vorlagen-Handling ---------------------------------------------- #

    @staticmethod
    def _remove_all_slides(prs):
        sld_id_lst = prs.slides._sldIdLst
        r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        for sld_id in list(sld_id_lst):
            prs.part.drop_rel(sld_id.get(r_ns))
            sld_id_lst.remove(sld_id)

    @staticmethod
    def _pick_content_layout(prs):
        """Layout mit Titel-Platzhalter und möglichst wenig weiteren Inhalts-
        Platzhaltern (z.B. «Nur Titel») – erhält Logo/Design der Vorlage."""
        layouts = _alle_layouts(prs)
        if not layouts:
            raise ValueError("Die Vorlage enthält keine Folienlayouts.")

        def score(lay):
            hat_titel = _titel_platzhalter(lay) is not None
            inhalt = sum(1 for p in lay.placeholders
                         if p.placeholder_format.type not in _TITEL_TYPEN
                         and p.placeholder_format.type not in _DEKOR_TYPEN)
            return (0 if hat_titel else 1, inhalt, len(lay.placeholders))

        return min(layouts, key=score)

    @staticmethod
    def _pick_title_layout(prs):
        layouts = _alle_layouts(prs)
        for lay in layouts:
            if any(p.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE
                   for p in lay.placeholders):
                return lay
        return layouts[0]

    def _add(self, layout):
        slide = self.prs.slides.add_slide(layout)
        self._clone_dekor(slide, layout)
        return slide

    @staticmethod
    def _clone_dekor(slide, layout):
        """Fusszeile/Datum/Seitenzahl vom Layout in die Folie übernehmen –
        python-pptx klont diese Platzhalter nicht automatisch."""
        vorhanden = {ph.placeholder_format.idx for ph in slide.placeholders}
        for ph in layout.placeholders:
            if (ph.placeholder_format.type in _DEKOR_TYPEN
                    and ph.placeholder_format.idx not in vorhanden):
                slide.shapes._spTree.append(deepcopy(ph._element))

    # ---- Grundgerüst --------------------------------------------------- #

    def _x(self, frac):
        return Emu(int(self.sw * frac))

    def _y(self, frac):
        return Emu(int(self.sh * frac))

    def _slide(self, titel):
        slide = self._add(self.content_layout)
        ph = _titel_platzhalter(slide)
        if ph is not None:
            ph.text = titel
            top = self._ph_unterkante(slide, ph)
        else:
            box = slide.shapes.add_textbox(self._x(0.06), self._y(0.05),
                                           self._x(0.70), self._y(0.12))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = titel
            p.font.size = Pt(26)
            p.font.bold = True
            top = 0.20
        return slide, top

    def _ph_unterkante(self, slide, ph):
        top, height = ph.top, ph.height
        if top is None or height is None:
            lay_ph = _titel_platzhalter(slide.slide_layout)
            if lay_ph is not None:
                top = lay_ph.top if top is None else top
                height = lay_ph.height if height is None else height
        if top is None or height is None:
            return 0.20
        frac = (int(top) + int(height)) / int(self.sh) + 0.02
        return min(max(frac, 0.14), 0.35)

    def _bullets(self, slide, items, top, size=16):
        box = slide.shapes.add_textbox(self._x(0.06), self._y(top),
                                       self._x(0.88), self._y(0.92 - top))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"•  {item}"
            p.font.size = Pt(size)
            p.space_after = Pt(8)
        return box

    def _style_tabelle(self, table, kopfzeile=True):
        """Einheitlich: schwarze Schrift überall, Kopfzeile hellgrau + fett –
        unabhängig vom (teils dunklen) Tabellenstil der Vorlage."""
        for r_i, row in enumerate(table.rows):
            for cell in row.cells:
                cell.text_frame.word_wrap = True
                if kopfzeile and r_i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _GRAU
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = _SCHWARZ
                    if kopfzeile and r_i == 0:
                        p.font.bold = True

    def _notizen(self, slide, text):
        slide.notes_slide.notes_text_frame.text = text

    # ---- Folien --------------------------------------------------------- #

    def titel(self, pia, datum):
        slide = self._add(self.title_layout)
        name = pia.get("projektname") or "Projekt"
        untertitel = "Projektinitialisierungsauftrag – Antrag auf Projektinitialisierungsfreigabe"
        meta = []
        if pia.get("projektleiter"):
            meta.append(f"Projektleitung: {pia['projektleiter']}")
        if pia.get("auftraggeber"):
            meta.append(f"Auftraggeber/in: {pia['auftraggeber']}")
        if datum:
            meta.append(datum)
        meta_text = "  ·  ".join(meta)

        titel_ph = _titel_platzhalter(slide)
        sub_ph = _platzhalter_nach_typ(slide, PP_PLACEHOLDER.SUBTITLE)
        body_ph = _platzhalter_nach_typ(slide, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        if titel_ph is not None:
            titel_ph.text = name
            if sub_ph is not None:
                sub_ph.text = untertitel
                if body_ph is not None and meta_text:
                    body_ph.text = meta_text
                elif meta_text:
                    sub_ph.text_frame.add_paragraph().text = meta_text
            elif body_ph is not None:
                body_ph.text = untertitel + ("\n" + meta_text if meta_text else "")
        else:
            box = slide.shapes.add_textbox(self._x(0.06), self._y(0.28),
                                           self._x(0.88), self._y(0.4))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(36)
            p.font.bold = True
            sub = tf.add_paragraph()
            sub.text = untertitel
            sub.font.size = Pt(18)
            if meta_text:
                m = tf.add_paragraph()
                m.text = meta_text
                m.font.size = Pt(14)
        self._notizen(slide, (
            "Begrüssen Sie den Auftraggeber und den Projektausschuss. "
            "Ziel der Sitzung: Freigabe des Projektinitialisierungsauftrags "
            "(Entscheid Projektinitialisierungsfreigabe). "
            "Kündigen Sie den Ablauf an: Ausgangslage, Ziele, Termine, Aufwand und "
            "Kosten, Risiken, Antrag."))

    def ausgangslage(self, pia, bullets):
        if not bullets:
            return
        slide, top = self._slide("Ausgangslage")
        self._bullets(slide, bullets, top)
        self._notizen(slide, (
            "Schildern Sie die Ausgangslage in eigenen Worten – nicht ablesen. "
            f"Kernbotschaft: {bullets[0]} "
            "Machen Sie den Handlungsbedarf deutlich: Warum jetzt, was passiert ohne das Projekt?"))

    def ziele(self, pia):
        ziele = pia.get("ziele") or []
        if not ziele:
            return
        slide, top = self._slide("Ziele der Phase Initialisierung")
        rows = ziele[:8]
        table = slide.shapes.add_table(
            len(rows) + 1, 2, self._x(0.06), self._y(top),
            self._x(0.88), self._y(min(0.1 * len(rows) + 0.05, 0.9 - top)),
        ).table
        table.columns[0].width = self._x(0.07)
        table.columns[1].width = self._x(0.81)
        table.cell(0, 0).text = "Nr."
        table.cell(0, 1).text = "Ziel"
        for r, ziel in enumerate(rows, 1):
            table.cell(r, 0).text = f"{r:02d}"
            table.cell(r, 1).text = ziel["beschreibung"]        # Volltext, keine Kürzung
            for p in table.cell(r, 1).text_frame.paragraphs:
                p.font.size = Pt(11)
        self._style_tabelle(table)
        self._notizen(slide, (
            "Gehen Sie die Ziele zügig durch, nicht wörtlich vorlesen. "
            "Kernaussage: Die Initialisierung schafft die Entscheidungsgrundlagen "
            "(insbesondere die Studie) – erst danach wird über die Umsetzung entschieden."))

    def termine(self, pia):
        termine = pia.get("termine") or []
        if not termine:
            return
        slide, top = self._slide("Ergebnisse und Termine")
        rows = termine[:10]
        table = slide.shapes.add_table(
            len(rows) + 1, 3, self._x(0.06), self._y(top),
            self._x(0.88), self._y(min(0.05 * (len(rows) + 1), 0.9 - top)),
        ).table
        table.columns[0].width = self._x(0.56)
        table.columns[1].width = self._x(0.14)
        table.columns[2].width = self._x(0.18)
        for c, kopf in enumerate(("Ergebnis", "Termin", "Abnahme durch")):
            table.cell(0, c).text = kopf
        for r, row in enumerate(rows, 1):
            werte = (row["ergebnis"], row.get("termin", ""), row.get("abnahme", ""))
            meilenstein = "meilenstein" in row["ergebnis"].lower()
            for c, wert in enumerate(werte):
                cell = table.cell(r, c)
                cell.text = wert
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.bold = meilenstein
        self._style_tabelle(table)
        letzter = max((t.get("termin", "") for t in termine), default="")
        self._notizen(slide, (
            "Erläutern Sie die HERMES-Logik der Reihenfolge: Die Analysen fliessen in "
            "die Studie; auf deren Basis wird der Entscheid 'Weiteres Vorgehen' gefällt; "
            "danach folgen Projektmanagementplan und Durchführungsauftrag. "
            f"Geplantes Ende der Initialisierung: {letzter or 'siehe Tabelle'}. "
            "Meilensteine sind fett hervorgehoben."))

    def gantt(self, pia):
        """Terminplan als gezeichnetes Gantt: Balken für Ergebnisse, Rauten
        (Dauer 0) für Meilensteine, Monatslinien als Orientierung."""
        eintraege = []
        for t in pia.get("termine") or []:
            datum = _parse_datum(t.get("termin"))
            if datum:
                eintraege.append((t["ergebnis"], datum,
                                  "meilenstein" in t["ergebnis"].lower()))
        if len(eintraege) < 3:
            return
        eintraege.sort(key=lambda e: e[1])
        eintraege = eintraege[:10]
        start = eintraege[0][1]
        ende = eintraege[-1][1]
        spann = max((ende - start).days, 1)

        slide, top = self._slide("Terminplan (Übersicht)")
        label_l, label_w = 0.06, 0.30
        plot_l, plot_r = label_l + label_w + 0.01, 0.94
        unten = 0.88
        zeile_h = min(0.065, (unten - top - 0.05) / len(eintraege))

        def x_von(d):
            return plot_l + (plot_r - plot_l) * ((d - start).days / spann)

        # Monatslinien + Beschriftung
        for m in _monatsanfaenge(start, ende):
            xf = x_von(m)
            linie = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, self._x(xf), self._y(top + 0.035),
                Emu(9525), self._y(unten - top - 0.035))
            linie.fill.solid()
            linie.fill.fore_color.rgb = _LINIENGRAU
            linie.line.fill.background()
            lbl = slide.shapes.add_textbox(self._x(xf - 0.03), self._y(top),
                                           self._x(0.06), self._y(0.035))
            p = lbl.text_frame.paragraphs[0]
            p.text = f"{_MONATE[m.month - 1]} {m.strftime('%y')}"
            p.font.size = Pt(9)
            p.font.color.rgb = _SCHWARZ

        vorher = start
        for i, (name, datum, meilenstein) in enumerate(eintraege):
            y = top + 0.05 + i * zeile_h
            lbl = slide.shapes.add_textbox(self._x(label_l), self._y(y),
                                           self._x(label_w), self._y(zeile_h))
            p = lbl.text_frame.paragraphs[0]
            lbl.text_frame.word_wrap = True
            p.text = _ohne_klammern(name)
            p.font.size = Pt(9)
            p.font.bold = meilenstein
            p.font.color.rgb = _SCHWARZ

            balken_h = zeile_h * 0.55
            if meilenstein:
                # Meilenstein: Dauer 0 – Raute am Termin.
                seite = self._y(balken_h * 1.4)
                raute = slide.shapes.add_shape(
                    MSO_SHAPE.DIAMOND,
                    Emu(int(self._x(x_von(datum))) - int(seite) // 2),
                    self._y(y + (zeile_h - balken_h * 1.4) / 2), seite, seite)
                raute.fill.solid()
                raute.fill.fore_color.rgb = _RAUTENDUNKEL
                raute.line.fill.background()
                ende_x = x_von(datum) + 0.006
            else:
                x1, x2 = x_von(vorher), x_von(datum)
                if x2 - x1 < 0.008:
                    x2 = x1 + 0.008
                balken = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, self._x(x1),
                    self._y(y + (zeile_h - balken_h) / 2),
                    self._x(x2 - x1), self._y(balken_h))
                balken.fill.solid()
                balken.fill.fore_color.rgb = _BALKENBLAU
                balken.line.fill.background()
                ende_x = x2 + 0.004
            # Das ECHTE Datum am Balkenende bzw. neben der Raute – die X-Achse
            # zeigt damit konkrete Termine, nicht Kalendertage ab Start.
            datum_lbl = slide.shapes.add_textbox(
                self._x(min(ende_x, 0.885)), self._y(y + zeile_h * 0.08),
                self._x(0.075), self._y(zeile_h * 0.8))
            dp = datum_lbl.text_frame.paragraphs[0]
            dp.text = datum.strftime("%d.%m.%y")
            dp.font.size = Pt(8)
            dp.font.bold = meilenstein
            dp.font.color.rgb = _SCHWARZ
            vorher = datum

        self._notizen(slide, (
            "Zeigen Sie den zeitlichen Ablauf: Balken entsprechen dem Zeitraum bis "
            "zum jeweiligen Termin, die Reihenfolge folgt den Abhängigkeiten. "
            "Die Rauten sind die Entscheid-Meilensteine (Dauer 0 Tage) – dort ist "
            "der Auftraggeber bzw. der Projektausschuss gefordert."))

    def personalaufwand(self, pia):
        personal = [p for p in (pia.get("personalaufwand") or []) if p.get("aufwand")]
        if not personal:
            return
        slide, top = self._slide("Personalaufwand der Initialisierung (PT)")
        data = CategoryChartData()
        data.categories = [
            p["rolle"] + (f"\n{p['name']}" if p.get("name") else "")
            for p in personal
        ]
        data.add_series("PT", [p["aufwand"] for p in personal])
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, self._x(0.06), self._y(top),
            self._x(0.88), self._y(min(0.68, 0.92 - top)), data,
        )
        chart = frame.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(11)
        try:
            chart.category_axis.tick_labels.font.size = Pt(10)
        except Exception:  # noqa: BLE001
            pass
        total = sum(p["aufwand"] for p in personal)
        groesste = max(personal, key=lambda p: p["aufwand"])
        self._notizen(slide, (
            f"Gesamtaufwand der Initialisierung: {total} Personentage. "
            f"Grösster Anteil: {groesste['rolle']} ({groesste['aufwand']} PT). "
            "Betonen Sie, dass dies NUR die Initialisierung betrifft – der Aufwand "
            "der Umsetzung wird erst nach der Studie geschätzt."))

    def kosten(self, pia):
        kosten = pia.get("kosten") or []
        einzel = [k for k in kosten
                  if k.get("betrag") and not _ist_summe(k["position"])]
        summen = [k for k in kosten if k.get("betrag") and _ist_summe(k["position"])]
        if not einzel and not summen:
            return
        slide, top = self._slide("Kosten der Initialisierung (CHF inkl. MwSt.)")
        if einzel:
            data = CategoryChartData()
            data.categories = [k["position"] for k in einzel]
            data.add_series("CHF", [k["betrag"] for k in einzel])
            frame = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, self._x(0.06), self._y(top),
                self._x(0.55), self._y(min(0.65, 0.92 - top)), data,
            )
            chart = frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
        if summen:
            box = slide.shapes.add_textbox(self._x(0.64), self._y(top + 0.05),
                                           self._x(0.32), self._y(0.5))
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for s in summen:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = f"{s['position']}: CHF {s['betrag']:,}".replace(",", "'")
                p.font.size = Pt(13)
                p.font.bold = "total" in s["position"].lower()
                p.space_after = Pt(6)
        total = next((s["betrag"] for s in summen if "total" in s["position"].lower()), None)
        extern = next((s["betrag"] for s in summen if "extern" in s["position"].lower()), None)
        hinweis = f"Gesamtkosten der Initialisierung: CHF {total:,}. ".replace(",", "'") if total else ""
        if extern:
            hinweis += f"Davon extern: CHF {extern:,}. ".replace(",", "'")
        self._notizen(slide, (
            hinweis + "Erläutern Sie kurz, wie die Kosten hergeleitet sind "
            "(Personalaufwand mal Tagessatz, plus Sachmittel). Auch hier gilt: "
            "nur die Phase Initialisierung, keine Umsetzungskosten."))

    def risiken(self, pia):
        risiken = pia.get("risiken") or []
        if not risiken:
            return
        platzierbar = [r for r in risiken if r.get("ew") and r.get("ag")]
        if platzierbar:
            self._risikomatrix(platzierbar, len(risiken))
        self._risiko_tabellen(risiken)

    def _risikomatrix(self, risiken, gesamt):
        slide, top = self._slide("Risikomatrix (Initialisierung)")
        stufen = ("Tief", "Mittel", "Hoch")
        table = slide.shapes.add_table(
            4, 4, self._x(0.14), self._y(top),
            self._x(0.72), self._y(min(0.6, 0.9 - top)),
        ).table
        table.cell(0, 0).text = "EW \\ AG"
        for c, s in enumerate(stufen, 1):
            table.cell(0, c).text = s
        for r, s in enumerate(reversed(stufen), 1):   # EW: hoch oben
            table.cell(r, 0).text = s
        for r in range(4):
            for c in range(4):
                cell = table.cell(r, c)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = _SCHWARZ
                if r == 0 or c == 0:
                    p.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _GRAU
                else:
                    ew = 4 - r
                    ag = c
                    score = ew * ag
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        _ROT if score >= 6 else _GELB if score >= 3 else _GRUEN)
        for risiko in risiken:
            r = 4 - risiko["ew"]
            c = risiko["ag"]
            cell = table.cell(r, c)
            p = cell.text_frame.paragraphs[0]
            bestehend = p.text
            neu = f"R{risiko['nr']}"
            p.text = f"{bestehend}  {neu}".strip() if bestehend else neu
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = _SCHWARZ
        rot = sum(1 for r in risiken if r["ew"] * r["ag"] >= 6)
        self._notizen(slide, (
            f"Lagebild: {gesamt} Risiken der Initialisierung, davon {rot} im roten "
            "Bereich. Botschaft: Die Risiken sind erkannt und mit Massnahmen hinterlegt "
            "– Details auf den Folgefolien. Gehen Sie hier nur auf die roten ein."))

    def _risiko_tabellen(self, risiken, pro_folie=4):
        """Risiken als Tabelle(n) mit VOLLEM Beschreibungstext (keine Kürzung)."""
        seiten = [risiken[i:i + pro_folie] for i in range(0, len(risiken), pro_folie)]
        for idx, seite in enumerate(seiten, 1):
            suffix = f" ({idx}/{len(seiten)})" if len(seiten) > 1 else ""
            slide, top = self._slide(f"Risiken der Initialisierung{suffix}")
            table = slide.shapes.add_table(
                len(seite) + 1, 4, self._x(0.06), self._y(top),
                self._x(0.88), self._y(min(0.16 * len(seite) + 0.05, 0.9 - top)),
            ).table
            table.columns[0].width = self._x(0.06)
            table.columns[1].width = self._x(0.62)
            table.columns[2].width = self._x(0.10)
            table.columns[3].width = self._x(0.10)
            for c, kopf in enumerate(("Nr.", "Risiko", "EW", "AG")):
                table.cell(0, c).text = kopf
            for r, risiko in enumerate(seite, 1):
                werte = (f"R{risiko['nr']}", risiko["beschreibung"],
                         _STUFEN_TEXT.get(risiko.get("ew"), ""),
                         _STUFEN_TEXT.get(risiko.get("ag"), ""))
                for c, wert in enumerate(werte):
                    cell = table.cell(r, c)
                    cell.text = wert
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
            self._style_tabelle(table)
            top_risiko = seite[0]
            self._notizen(slide, (
                "Nicht alle Risiken vorlesen – heben Sie das gewichtigste hervor: "
                f"R{top_risiko['nr']}: {_erster_satz(top_risiko['beschreibung'])} "
                "Botschaft: Zu jedem Risiko existiert eine Massnahme (im PIA, Kapitel Risiken)."))

    def antrag(self, pia):
        slide, top = self._slide("Antrag an den Auftraggeber")
        items = [
            "Der Projektinitialisierungsauftrag liegt in der vorliegenden Fassung zur "
            "Freigabe vor.",
            "Antrag: Erteilung der Projektinitialisierungsfreigabe und Start der "
            "Phase Initialisierung.",
        ]
        naechste = [t for t in (pia.get("termine") or [])
                    if "meilenstein" in t.get("ergebnis", "").lower()]
        for t in naechste[:2]:
            zeile = _ohne_klammern(t["ergebnis"])
            if t.get("termin"):
                zeile += f" – geplant {t['termin']}"
            items.append(zeile)
        self._bullets(slide, items, top, size=16)
        self._notizen(slide, (
            "Bitten Sie den Auftraggeber EXPLIZIT um den Entscheid: Erteilung der "
            "Projektinitialisierungsfreigabe. Klären Sie offene Fragen jetzt; halten "
            "Sie den Entscheid und allfällige Auflagen im Protokoll (Liste "
            "Projektentscheide Steuerung) fest."))

    def schluss(self, pia, datum):
        if self.closing_layout is None:
            return
        slide = self._add(self.closing_layout)
        titel_ph = _titel_platzhalter(slide)
        sub_ph = _platzhalter_nach_typ(slide, PP_PLACEHOLDER.SUBTITLE,
                                       PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        zeilen = [pia.get("projektname") or ""]
        if pia.get("projektleiter"):
            zeilen.append(f"Projektleitung: {pia['projektleiter']}")
        if datum:
            zeilen.append(datum)
        if titel_ph is not None:
            titel_ph.text = "Besten Dank."
            if sub_ph is not None:
                sub_ph.text = "\n".join(z for z in zeilen if z)
        else:
            box = slide.shapes.add_textbox(self._x(0.06), self._y(0.35),
                                           self._x(0.88), self._y(0.3))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "Besten Dank."
            p.font.size = Pt(32)
            p.font.bold = True
        self._notizen(slide, (
            "Bedanken Sie sich für die Aufmerksamkeit. Nennen Sie den nächsten "
            "Schritt: Nach erteilter Freigabe startet die Phase Initialisierung "
            "gemäss Terminplan."))
